from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation
prs = Presentation()

# Slides content
slides = [
    ("STAC Catalog Best Practices & Asset Guidelines", "Prepared by: Vishnu S\nDate: 01-Sep-2025"),
    ("Agenda", "1. STAC Metadata Best Practices\n2. Common Metadata Fields\n3. Assets and Roles\n4. Legend JSON in STAC\n5. Remote Data Download Assets\n6. References & Example Catalogs"),
    ("STAC Metadata Best Practices", "- Follow STAC core specification for interoperability.\n- Use extensions for domain-specific metadata (eo, raster, classification).\n- Maintain consistent naming conventions for fields.\n- Include dates, geometry, and platform info for context."),
    ("Common Metadata Fields", "stac_version, stac_extensions, datetime (optional), start_datetime / end_datetime, geometry / bbox, license, providers (Producer, Licensor, Host), platform, constellation, GSD, first_acquisition / last_acquisition, acquired"),
    ("Assets in STAC Items", "Each item can include multiple assets with href, title, description, type (media type), roles (thumbnail, overview, data, metadata)"),
    ("Legend JSON in STAC", "- Optional but useful for visual interpretation.\n- Used in classified datasets (land cover maps).\n- Added as an asset in the STAC item.\nExample: {\"href\": \"https://example.com/legend.json\", \"title\": \"Classification Legend\", \"description\": \"Maps class IDs to color codes.\"}"),
    ("Remote Data Download Assets", "- Provide direct access to datasets.\n- Assign unique href per asset.\n- Include title, description, media type.\n- Use roles: data, thumbnail, metadata.\nExample: {\"thumbnail\": {\"href\": \"https://example.com/thumbnail.jpg\", \"title\": \"Thumbnail Image\", \"media_type\": \"image/jpeg\"}, \"data\": {\"href\": \"https://example.com/data.zip\", \"title\": \"Full Dataset\", \"media_type\": \"application/zip\"}}"),
    ("Example STAC Catalogs", "1. Planet Labs STAC Catalog\n2. OpenLandMap STAC\n3. Microsoft Planetary Computer STAC\n4. Radiant Earth STAC Examples\n5. Additional catalogs on STAC Index"),
    ("Key Takeaways", "- Use common metadata fields consistently.\n- Legend JSON assets improve interpretability.\n- Remote download assets enhance usability.\n- Reference existing STAC catalogs for best practices.\n- Extensions capture domain-specific requirements."),
    ("References", "- https://github.com/radiantearth/stac-spec/blob/v1.1.0/commons/common-metadata.md\n- Planet Labs STAC Catalog\n- OpenLandMap STAC\n- Microsoft Planetary Computer STAC\n- ODC-STAC Best Practices\n- Pystac Asset Documentation")
]

# Add slides
for title, content in slides:
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save presentation
output_path = "/mnt/data/STAC_Best_Practices.pptx"
prs.save(output_path)
output_path
